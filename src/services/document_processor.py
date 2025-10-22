"""Document processing service for text extraction and chunking."""

import asyncio
import hashlib
import io
import uuid
from pathlib import Path
from typing import Dict, List, Optional, Any, Union
from datetime import datetime
import mimetypes

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, update

from ..core.database import get_db_session
from ..core.logging import get_logger, LoggerMixin
from ..models.document import Document, DocumentStatus
from ..models.chunk import Chunk
from ..services.embedding_service import embedding_service
from ..services.storage import storage_service, StorageService

logger = get_logger(__name__)


class DocumentProcessor(LoggerMixin):
    """Service for processing uploaded documents."""
    
    def __init__(self, storage_service: StorageService):
        self.storage_service = storage_service
        self.chunk_size = 1000  # Default chunk size in characters
        self.chunk_overlap = 200  # Overlap between chunks
        
    async def process_document_async(self, document_id: str) -> bool:
        """Process document asynchronously."""
        try:
            async with get_db_session() as session:
                # Get document
                document = await session.get(Document, document_id)
                if not document:
                    self.log_error("Document not found", document_id=document_id)
                    return False
                
                # Update status to processing
                document.status = DocumentStatus.PROCESSING
                await session.commit()
                
                # Download file content
                file_content = await self.storage_service.get_file(document.file_path)
                if not file_content:
                    await self._mark_document_failed(session, document, "File not found in storage")
                    return False
                
                # Extract text based on file type
                text_content = await self._extract_text(file_content, document.file_type, document.filename)
                
                if not text_content:
                    await self._mark_document_failed(session, document, "No text content extracted")
                    return False
                
                # Create chunks
                chunks = self._create_chunks(text_content, document_id)
                document.total_chunks = len(chunks)
                
                # Store chunks in database
                chunk_objects = []
                for i, chunk_content in enumerate(chunks):
                    chunk = Chunk(
                        id=str(uuid.uuid4()),
                        document_id=document_id,
                        content=chunk_content,
                        chunk_index=i,
                        start_offset=None,  # Could be calculated if needed
                        end_offset=None,
                        metadata={
                            "word_count": len(chunk_content.split()),
                            "char_count": len(chunk_content)
                        }
                    )
                    chunk_objects.append(chunk)
                
                session.add_all(chunk_objects)
                
                # Generate embeddings for chunks
                if chunks:
                    try:
                        chunk_texts = [chunk.content for chunk in chunk_objects]
                        embeddings = await embedding_service.generate_embeddings(
                            chunk_texts,
                            use_cache=True
                        )
                        
                        # Update chunks with embedding IDs (simplified - in production you'd store in vector DB)
                        for chunk, embedding in zip(chunk_objects, embeddings):
                            chunk.embedding_id = str(uuid.uuid4())  # Placeholder
                            
                    except Exception as e:
                        self.log_warning("Embedding generation failed", document_id=document_id, error=e)
                
                # Update document status
                document.status = DocumentStatus.PROCESSED
                document.processed_chunks = len(chunks)
                document.last_updated = datetime.utcnow()
                
                await session.commit()
                
                self.log_info(
                    "Document processed successfully",
                    document_id=document_id,
                    chunks_created=len(chunks)
                )
                
                return True
                
        except Exception as e:
            self.log_error("Document processing failed", document_id=document_id, error=e)
            async with get_db_session() as session:
                document = await session.get(Document, document_id)
                if document:
                    await self._mark_document_failed(session, document, str(e))
            return False
    
    async def _extract_text(self, file_content: bytes, file_type: str, filename: str) -> str:
        """Extract text from file content based on file type."""
        try:
            # Get file extension
            file_path = Path(filename)
            extension = file_path.suffix.lower()
            
            # Text files
            if extension in ['.txt', '.md', '.py', '.js', '.html', '.htm', '.css', '.json', '.csv']:
                return self._extract_text_file(file_content)
            
            # PDF files
            elif extension == '.pdf':
                return await self._extract_pdf_text(file_content)
            
            # Word documents
            elif extension in ['.docx', '.doc']:
                return await self._extract_word_text(file_content)
            
            # Excel files
            elif extension in ['.xlsx', '.xls']:
                return await self._extract_excel_text(file_content)
            
            # PowerPoint files
            elif extension in ['.pptx', '.ppt']:
                return await self._extract_powerpoint_text(file_content)
            
            # Fallback: try to decode as text
            else:
                return self._extract_text_file(file_content)
                
        except Exception as e:
            self.log_error("Text extraction failed", filename=filename, error=e)
            raise
    
    def _extract_text_file(self, file_content: bytes) -> str:
        """Extract text from text-based files."""
        # Try different encodings
        for encoding in ['utf-8', 'latin-1', 'cp1252']:
            try:
                return file_content.decode(encoding)
            except UnicodeDecodeError:
                continue
        
        # Fallback: decode with errors ignored
        return file_content.decode('utf-8', errors='ignore')
    
    async def _extract_pdf_text(self, file_content: bytes) -> str:
        """Extract text from PDF files."""
        try:
            import PyPDF2
            
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = []
            for page in pdf_reader.pages:
                text_content.append(page.extract_text())
            
            return '\n'.join(text_content)
            
        except ImportError:
            # Fallback if PyPDF2 is not available
            self.log_warning("PyPDF2 not available, using text extraction fallback")
            return self._extract_text_file(file_content)
        except Exception as e:
            self.log_error("PDF text extraction failed", error=e)
            raise
    
    async def _extract_word_text(self, file_content: bytes) -> str:
        """Extract text from Word documents."""
        try:
            import docx
            
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            
            text_content = []
            for paragraph in doc.paragraphs:
                text_content.append(paragraph.text)
            
            return '\n'.join(text_content)
            
        except ImportError:
            self.log_warning("python-docx not available, using text extraction fallback")
            return self._extract_text_file(file_content)
        except Exception as e:
            self.log_error("Word text extraction failed", error=e)
            raise
    
    async def _extract_excel_text(self, file_content: bytes) -> str:
        """Extract text from Excel files."""
        try:
            import pandas as pd
            
            excel_file = io.BytesIO(file_content)
            
            # Read all sheets
            excel_data = pd.read_excel(excel_file, sheet_name=None)
            
            text_content = []
            for sheet_name, df in excel_data.items():
                text_content.append(f"Sheet: {sheet_name}")
                text_content.append(df.to_string(index=False))
                text_content.append("")
            
            return '\n'.join(text_content)
            
        except ImportError:
            self.log_warning("pandas not available for Excel processing")
            return self._extract_text_file(file_content)
        except Exception as e:
            self.log_error("Excel text extraction failed", error=e)
            raise
    
    async def _extract_powerpoint_text(self, file_content: bytes) -> str:
        """Extract text from PowerPoint files."""
        try:
            import pptx
            
            ppt_file = io.BytesIO(file_content)
            presentation = pptx.Presentation(ppt_file)
            
            text_content = []
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_content.append(f"Slide {slide_num}:")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                
                text_content.append("")
            
            return '\n'.join(text_content)
            
        except ImportError:
            self.log_warning("python-pptx not available for PowerPoint processing")
            return self._extract_text_file(file_content)
        except Exception as e:
            self.log_error("PowerPoint text extraction failed", error=e)
            raise
    
    def _create_chunks(self, text: str, document_id: str) -> List[str]:
        """Create chunks from text content."""
        if not text:
            return []
        
        # Simple sentence-aware chunking
        sentences = self._split_into_sentences(text)
        chunks = []
        current_chunk = ""
        
        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk) + len(sentence) + 1 > self.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    
                    # Start new chunk with overlap
                    if self.chunk_overlap > 0 and len(current_chunk) > self.chunk_overlap:
                        current_chunk = current_chunk[-self.chunk_overlap:] + " " + sentence
                    else:
                        current_chunk = sentence
                else:
                    # Sentence is longer than chunk size, split it
                    if len(sentence) > self.chunk_size:
                        word_chunks = self._split_long_sentence(sentence)
                        chunks.extend(word_chunks[:-1])
                        current_chunk = word_chunks[-1] if word_chunks else ""
                    else:
                        current_chunk = sentence
            else:
                current_chunk += " " + sentence if current_chunk else sentence
        
        # Add the last chunk
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences."""
        import re
        
        # Simple sentence splitting (can be improved with NLTK or spaCy)
        sentence_endings = re.compile(r'[.!?]+')
        sentences = sentence_endings.split(text)
        
        # Clean and filter sentences
        sentences = [s.strip() for s in sentences if s.strip()]
        
        return sentences
    
    def _split_long_sentence(self, sentence: str) -> List[str]:
        """Split a long sentence into smaller chunks."""
        words = sentence.split()
        chunks = []
        current_chunk = []
        current_length = 0
        
        for word in words:
            if current_length + len(word) + 1 > self.chunk_size:
                if current_chunk:
                    chunks.append(' '.join(current_chunk))
                    current_chunk = [word]
                    current_length = len(word)
                else:
                    # Word is longer than chunk size, just add it
                    chunks.append(word)
                    current_chunk = []
                    current_length = 0
            else:
                current_chunk.append(word)
                current_length += len(word) + 1
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    async def _mark_document_failed(self, session: AsyncSession, document: Document, error_message: str):
        """Mark document as failed with error message."""
        document.status = DocumentStatus.FAILED
        document.error_message = error_message
        document.last_updated = datetime.utcnow()
        await session.commit()
        
        self.log_error(
            "Document marked as failed",
            document_id=document.id,
            error_message=error_message
        )
    
    async def reprocess_document(self, document_id: str) -> bool:
        """Reprocess an existing document."""
        try:
            async with get_db_session() as session:
                # Delete existing chunks
                await session.execute(
                    select(Chunk).where(Chunk.document_id == document_id)
                )
                result = await session.execute(
                    select(Chunk).where(Chunk.document_id == document_id)
                )
                existing_chunks = result.scalars().all()
                
                for chunk in existing_chunks:
                    await session.delete(chunk)
                
                await session.commit()
            
            # Process document again
            return await self.process_document_async(document_id)
            
        except Exception as e:
            self.log_error("Document reprocessing failed", document_id=document_id, error=e)
            return False
    
    async def get_processing_status(self, document_id: str) -> Dict[str, Any]:
        """Get document processing status."""
        try:
            async with get_db_session() as session:
                document = await session.get(Document, document_id)
                
                if not document:
                    return {"status": "not_found"}
                
                return {
                    "status": document.status.value,
                    "processed_chunks": document.processed_chunks,
                    "total_chunks": document.total_chunks,
                    "progress": (
                        (document.processed_chunks / document.total_chunks * 100)
                        if document.total_chunks > 0 else 0
                    ),
                    "error_message": document.error_message,
                    "last_updated": document.last_updated
                }
                
        except Exception as e:
            self.log_error("Failed to get processing status", document_id=document_id, error=e)
            return {"status": "error", "error": str(e)}
    
    async def health_check(self) -> Dict[str, Any]:
        """Check document processor health."""
        return {
            "status": "healthy",
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
            "supported_formats": [".txt", ".pdf", ".docx", ".xlsx", ".pptx", ".md", ".html", ".json", ".csv"]
        }